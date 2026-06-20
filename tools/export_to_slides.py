#!/usr/bin/env python3
"""
ThreadIntel — Export research report to PowerPoint (.pptx) or Google Slides.

Professional-grade slide deck generator. Each deck includes:
  - Branded title slide with date and metadata
  - Executive summary with AI-generated insights
  - Sentiment analysis bar chart
  - Top findings with source badges and sentiment indicators
  - Source breakdown chart
  - Comparison slides (auto-detected for "X vs Y" topics)
  - Branded footer on every slide

Usage:
    python tools/export_to_slides.py "AI tools 2026"
    python tools/export_to_slides.py "Notion vs Linear" --layout comparison
    python tools/export_to_slides.py "SaaS trends" --layout detailed
    python tools/export_to_slides.py "AI tools 2026" --slides  # → Google Slides URL

Layouts:
    auto        picks best layout for the topic type (default)
    overview    title + summary + top findings grid (most shareable)
    detailed    one findings slide per source group
    comparison  side-by-side topic analysis (auto for "X vs Y")
    bullets     tight bullet slides, good for internal review
"""

import os, sys, re
from pathlib import Path
from datetime import datetime
from collections import Counter

ROOT = Path(__file__).parent.parent
sys.path.insert(0, str(ROOT / "threadintel"))
from dotenv import load_dotenv
load_dotenv(ROOT / ".env")

CREDS_FILE = ROOT / "credentials.json"

# ─── Palette ──────────────────────────────────────────────────────────────────
BG          = (10,  12,  26)    # deep navy
BG2         = (18,  20,  42)    # slightly lighter panel bg
ACCENT      = (99, 157, 255)    # electric blue
ACCENT2     = (147, 112, 255)   # purple
WHITE       = (255, 255, 255)
OFF_WHITE   = (220, 225, 240)
MUTED       = (140, 148, 175)
GREEN       = (72,  199, 142)
RED         = (255,  90,  90)
YELLOW      = (255, 196,  57)
DIVIDER     = (30,  35,  65)

# ─── pptx helpers ─────────────────────────────────────────────────────────────

def _rgb(r, g, b):
    from pptx.dml.color import RGBColor
    return RGBColor(r, g, b)


def _set_slide_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(*color)


def _rect(slide, l, t, w, h, fill_color, alpha=None):
    from pptx.util import Inches, Pt
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h),
    )
    shape.line.fill.background()  # no border
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(*fill_color)
    else:
        shape.fill.background()
    return shape


def _label(slide, l, t, w, h, text,
           size=12, bold=False, color=WHITE,
           align="left", wrap=True, italic=False):
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf  = box.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = {
        "left":   PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right":  PP_ALIGN.RIGHT,
    }.get(align, PP_ALIGN.LEFT)

    run = p.add_run()
    run.text           = str(text)
    run.font.size      = Pt(size)
    run.font.bold      = bold
    run.font.italic    = italic
    run.font.color.rgb = _rgb(*color)
    return box


def _bar(slide, l, t, w, h, pct, fill_color, bg_color=DIVIDER):
    """Draw a simple progress bar."""
    _rect(slide, l, t, w, h, bg_color)
    bar_w = max(w * min(pct / 100, 1), 0.02)
    _rect(slide, l, t, bar_w, h, fill_color)


def _accent_line(slide, l, t, w, color=ACCENT):
    """Thin horizontal accent rule."""
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.dml import MSO_THEME_COLOR
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(0.025))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(*color)
    shape.line.fill.background()


def _source_badge_color(source: str):
    src = source.lower()
    if "reddit"    in src: return (255, 87, 34)
    if "hacker"    in src: return (255, 102, 0)
    if "news"      in src: return (33, 150, 243)
    if "duck"      in src or "web" in src: return (76, 175, 80)
    if "stack"     in src: return (244, 130, 52)
    if "product"   in src: return (218, 81, 61)
    return (99, 157, 255)


def _sent_color(s: str):
    s = (s or "").lower()
    if s == "positive": return GREEN
    if s == "negative": return RED
    return YELLOW


# ─── Slide 1: Title ───────────────────────────────────────────────────────────

def _slide_title(prs, topic, date_str, n_findings, sources_list):
    from pptx.util import Inches
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)

    # Left colour band
    _rect(slide, 0, 0, 0.07, 5.625, ACCENT)

    # Eyebrow
    _label(slide, 0.35, 0.55, 9.3, 0.4,
           "THREADINTEL  ·  MARKET RESEARCH REPORT",
           size=9, color=MUTED, bold=True)

    # Main title — two lines if long
    words = topic.split()
    if len(words) > 6:
        mid = len(words) // 2
        line1 = " ".join(words[:mid])
        line2 = " ".join(words[mid:])
        _label(slide, 0.35, 1.1,  9.3, 1.0, line1, size=34, bold=True, color=WHITE)
        _label(slide, 0.35, 1.95, 9.3, 1.0, line2, size=34, bold=True, color=WHITE)
        meta_y = 3.2
    else:
        _label(slide, 0.35, 1.2, 9.3, 1.4, topic, size=38, bold=True, color=WHITE)
        meta_y = 2.9

    _accent_line(slide, 0.35, meta_y - 0.1, 2.2, ACCENT)

    # Meta row
    _label(slide, 0.35, meta_y,       2.5, 0.35, date_str,          size=11, color=OFF_WHITE)
    _label(slide, 0.35, meta_y + 0.4, 2.5, 0.35, f"{n_findings} data points", size=11, color=MUTED)
    _label(slide, 0.35, meta_y + 0.8, 7.0, 0.35,
           "Sources: " + " · ".join(sources_list[:6]), size=10, color=MUTED)

    # Powered-by tag bottom right
    _label(slide, 7.5, 5.1, 2.2, 0.3, "threadintel.io",
           size=9, color=MUTED, align="right", italic=True)


# ─── Slide 2: Executive Summary ───────────────────────────────────────────────

def _slide_summary(prs, summary, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _accent_line(slide, 0.4, 0.3, 1.4, ACCENT)

    _label(slide, 0.4, 0.45, 9.0, 0.55, "Executive Summary",
           size=22, bold=True, color=WHITE)
    _rect(slide, 0.4, 1.1, 9.2, 0.02, DIVIDER)

    # Summary paragraph
    summary_text = (summary or "")[:380]
    _label(slide, 0.4, 1.2, 9.0, 1.3, summary_text, size=13, color=OFF_WHITE)

    # Bullet takeaways
    if bullets:
        _label(slide, 0.4, 2.65, 4.0, 0.38, "KEY TAKEAWAYS",
               size=9, bold=True, color=MUTED)
        y = 3.1
        for b in bullets[:5]:
            _rect(slide, 0.4, y + 0.07, 0.18, 0.18, ACCENT)
            _label(slide, 0.72, y, 8.7, 0.45, str(b)[:120], size=11.5, color=OFF_WHITE)
            y += 0.55

    _label(slide, 7.5, 5.2, 2.2, 0.25, "threadintel.io",
           size=8, color=MUTED, align="right", italic=True)


# ─── Slide 3: Sentiment bar chart ─────────────────────────────────────────────

def _slide_sentiment(prs, findings):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _accent_line(slide, 0.4, 0.3, 1.4, ACCENT)
    _label(slide, 0.4, 0.45, 9.0, 0.55, "Sentiment Analysis",
           size=22, bold=True, color=WHITE)
    _rect(slide, 0.4, 1.1, 9.2, 0.02, DIVIDER)

    n   = max(len(findings), 1)
    pos = round(sum(1 for f in findings if f.get("sentiment") == "positive") / n * 100)
    neg = round(sum(1 for f in findings if f.get("sentiment") == "negative") / n * 100)
    neu = 100 - pos - neg

    bars = [
        ("Positive Signals", pos, GREEN,  "Constructive mentions, recommendations, satisfaction"),
        ("Neutral / Mixed",  neu, YELLOW, "Informational, balanced, or ambiguous discussions"),
        ("Negative Signals", neg, RED,    "Complaints, frustrations, criticism"),
    ]

    y = 1.5
    for label, pct, color, desc in bars:
        _label(slide, 0.5, y,        2.2, 0.35, label, size=12, bold=True, color=color)
        _label(slide, 0.5, y + 0.38, 2.2, 0.30, desc,  size=9,  color=MUTED)
        _bar(slide,   2.8, y + 0.05, 5.8, 0.35, pct, color)
        _label(slide, 8.75, y + 0.05, 0.9, 0.35, f"{pct}%",
               size=14, bold=True, color=color, align="right")
        y += 1.1

    _label(slide, 7.5, 5.2, 2.2, 0.25, "threadintel.io",
           size=8, color=MUTED, align="right", italic=True)


# ─── Slide 4: Findings cards (2-column grid) ──────────────────────────────────

def _slide_findings(prs, items, heading="Top Findings", sub=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _accent_line(slide, 0.4, 0.3, 1.4, ACCENT)
    _label(slide, 0.4, 0.45, 9.0, 0.55, heading, size=22, bold=True, color=WHITE)
    if sub:
        _label(slide, 0.4, 0.95, 9.0, 0.3, sub, size=10, color=MUTED, italic=True)
    _rect(slide, 0.4, 1.1, 9.2, 0.02, DIVIDER)

    grid = items[:6]
    positions = [
        (0.28, 1.22), (5.12, 1.22),
        (0.28, 2.72), (5.12, 2.72),
        (0.28, 4.22), (5.12, 4.22),
    ]

    for i, f in enumerate(grid):
        x, y = positions[i]
        src   = f.get("source", "")
        title = f.get("title", "")[:75]
        text  = (f.get("text") or "")[:110].strip()
        sent  = f.get("sentiment", "neutral")
        scol  = _sent_color(sent)
        bcol  = _source_badge_color(src)

        # Card background
        _rect(slide, x, y, 4.6, 1.38, BG2)
        # Source badge
        _rect(slide, x, y, 1.35, 0.22, bcol)
        _label(slide, x + 0.05, y, 1.25, 0.22,
               src.replace("Web (DuckDuckGo)", "Web").upper(), size=7.5, bold=True, color=WHITE)
        # Sentiment dot
        _rect(slide, x + 4.35, y + 0.02, 0.18, 0.18, scol)
        # Title
        _label(slide, x + 0.1, y + 0.27, 4.35, 0.52, title,
               size=10.5, bold=True, color=WHITE, wrap=True)
        # Snippet
        _label(slide, x + 0.1, y + 0.83, 4.35, 0.48, text,
               size=8.5, color=MUTED, wrap=True)

    _label(slide, 7.5, 5.2, 2.2, 0.25, "threadintel.io",
           size=8, color=MUTED, align="right", italic=True)


# ─── End slide: ThreadIntel branding ──────────────────────────────────────────

def _slide_end(prs, topic: str):
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)

    # Gradient accent line across centre
    _rect(slide, 1.0, 2.6, 8.0, 0.04, ACCENT)

    # Logo text
    tf = slide.shapes.add_textbox(
        _i(3.5), _i(1.5), _i(3.0), _i(0.8)
    ).text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r1 = p.add_run(); r1.text = "Thread"
    r1.font.size = Pt(32); r1.font.bold = True
    r1.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    r2 = p.add_run(); r2.text = "Intel"
    r2.font.size = Pt(32); r2.font.bold = True
    r2.font.color.rgb = RGBColor(0x7C, 0x6F, 0xFF)

    _label(slide, 0.5, 2.85, 9.0, 0.4,
           "AI-powered market intelligence from Reddit, Hacker News, and the open web",
           size=11, color=MUTED, align="center")
    _label(slide, 0.5, 3.4, 9.0, 0.4,
           "threadintel.io  ·  $9.99/month  ·  Unlimited reports",
           size=10, color=MUTED, align="center")
    _label(slide, 0.5, 4.0, 9.0, 0.4,
           f"Report: {topic[:70]}",
           size=9, color=RGBColor(0x4A, 0x55, 0x68), align="center", italic=True)


# ─── Slide 5: Sources breakdown ───────────────────────────────────────────────

def _slide_sources(prs, findings):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _accent_line(slide, 0.4, 0.3, 1.4, ACCENT)
    _label(slide, 0.4, 0.45, 9.0, 0.55, "Sources Breakdown",
           size=22, bold=True, color=WHITE)
    _rect(slide, 0.4, 1.1, 9.2, 0.02, DIVIDER)

    counts = Counter(f.get("source", "Unknown") for f in findings)
    total  = max(sum(counts.values()), 1)
    top    = counts.most_common(7)

    y = 1.35
    for src, cnt in top:
        pct  = round(cnt / total * 100)
        bcol = _source_badge_color(src)
        src_label = src.replace("Web (DuckDuckGo)", "DuckDuckGo")[:28]
        _rect(slide, 0.45, y + 0.06, 0.18, 0.18, bcol)
        _label(slide, 0.75, y, 2.8, 0.38, src_label, size=12, bold=True, color=OFF_WHITE)
        _bar(slide, 3.6, y + 0.06, 5.0, 0.27, pct, bcol)
        _label(slide, 8.72, y, 1.0, 0.38,
               f"{cnt}", size=12, bold=True, color=bcol, align="right")
        _label(slide, 8.72, y + 0.35, 1.0, 0.22,
               f"{pct}%", size=9, color=MUTED, align="right")
        y += 0.62

    _label(slide, 7.5, 5.2, 2.2, 0.25, "threadintel.io",
           size=8, color=MUTED, align="right", italic=True)


# ─── Slide 6: Comparison (X vs Y) ─────────────────────────────────────────────

def _slide_comparison(prs, findings, sides, heading="Comparison"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide)
    _accent_line(slide, 0.4, 0.3, 1.4, ACCENT)
    _label(slide, 0.4, 0.45, 9.0, 0.55, heading,
           size=22, bold=True, color=WHITE)
    _rect(slide, 0.4, 1.1, 9.2, 0.02, DIVIDER)

    sides = sides[:3]
    col_w = 9.2 / len(sides)
    colors = [ACCENT, ACCENT2, GREEN]

    for i, side in enumerate(sides):
        x = 0.3 + i * col_w
        c = colors[i % len(colors)]
        side_items = [f for f in findings
                      if side.lower() in (f.get("title") or "").lower()
                      or side.lower() in (f.get("text") or "").lower()][:5]

        # Column header
        _rect(slide, x, 1.22, col_w - 0.15, 0.45, c)
        _label(slide, x + 0.08, 1.27, col_w - 0.3, 0.35,
               side.upper(), size=12, bold=True, color=BG, align="center")

        # Bullet items
        y_item = 1.82
        for f in side_items:
            title = (f.get("title") or "")[:58]
            _rect(slide, x + 0.08, y_item + 0.09, 0.1, 0.1, c)
            _label(slide, x + 0.25, y_item, col_w - 0.4, 0.48,
                   title, size=9, color=OFF_WHITE, wrap=True)
            y_item += 0.58

        if not side_items:
            _label(slide, x + 0.1, 2.0, col_w - 0.2, 0.4,
                   "No direct mentions found", size=10, color=MUTED, italic=True)

    _label(slide, 7.5, 5.2, 2.2, 0.25, "threadintel.io",
           size=8, color=MUTED, align="right", italic=True)


# ─── Master builder ───────────────────────────────────────────────────────────

def export_to_pptx(
    topic: str,
    findings=None,
    sources: list = None,
    layout: str = "auto",
    output_path: str = None,
) -> str:
    from pptx import Presentation
    from pptx.util import Inches
    from research import smart_research
    from email_brief import generate_content, detect_report_type

    # ── Research ──────────────────────────────────────────────────────────────
    if findings is None:
        print(f"  Researching: {topic}")
        findings = smart_research(topic, sources=sources)

    # generate_content expects the structured dict; keep original for Groq prompt
    findings_raw = findings

    # ── Report content ────────────────────────────────────────────────────────
    print(f"  Generating report content...")
    report_data = generate_content(topic, findings=findings_raw)

    # Flatten for slide building (list of finding dicts)
    if isinstance(findings, dict):
        flat = []
        for key in ("reddit","hackernews","news","web","stackoverflow","producthunt","newsapi","indiehackers","github_issues"):
            flat.extend(findings.get(key, []))
        findings = flat
    findings = [f for f in findings if isinstance(f, dict)]
    report_type  = detect_report_type(topic)

    summary = str(report_data.get("summary") or "")
    bullets = [str(b) for b in (report_data.get("bullets") or report_data.get("tldr") or []) if b]

    # ── Layout selection ──────────────────────────────────────────────────────
    if layout == "auto":
        layout = "comparison" if report_type == "comparison" else "overview"

    # ── Build deck ────────────────────────────────────────────────────────────
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(5.625)   # 16:9

    n      = datetime.now()
    date_s = f"{n.strftime('%B')} {n.day}, {n.year}"
    srcs   = sorted(set(
        f.get("source","").replace("Web (DuckDuckGo)","DuckDuckGo")
        for f in findings if f.get("source")
    ))

    # Slide 1 — Title
    _slide_title(prs, topic, date_s, len(findings), srcs)

    # Slide 2 — Executive Summary
    _slide_summary(prs, summary, bullets)

    # Slide 3 — Sentiment
    _slide_sentiment(prs, findings)

    # Slides 4+ — Findings
    if layout == "comparison":
        sides = re.split(r"\s+vs\.?\s+", topic, flags=re.IGNORECASE)
        if len(sides) > 1:
            _slide_comparison(prs, findings, sides)
        _slide_findings(prs, findings[:6], "Supporting Evidence")
        if len(findings) > 6:
            _slide_findings(prs, findings[6:12], "More Findings")

    elif layout == "detailed":
        by_src: dict = {}
        for f in findings:
            by_src.setdefault(f.get("source","Other"), []).append(f)
        for src_name, items in by_src.items():
            for chunk_start in range(0, len(items), 6):
                _slide_findings(prs, items[chunk_start:chunk_start+6],
                                f"{src_name} — Findings",
                                f"{len(items)} total results from this source")

    elif layout == "bullets":
        chunk_size = 10
        for i in range(0, min(len(findings), 50), chunk_size):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            _set_slide_bg(slide)
            _accent_line(slide, 0.4, 0.3, 1.4, ACCENT)
            _label(slide, 0.4, 0.45, 9.0, 0.55,
                   f"All Findings — {i+1} to {min(i+chunk_size, len(findings))}",
                   size=20, bold=True, color=WHITE)
            _rect(slide, 0.4, 1.1, 9.2, 0.02, DIVIDER)
            y = 1.25
            for f in findings[i:i+chunk_size]:
                src   = f.get("source","")
                title = (f.get("title") or "")[:90]
                bcol  = _source_badge_color(src)
                scol  = _sent_color(f.get("sentiment",""))
                _rect(slide, 0.4, y + 0.07, 0.12, 0.12, bcol)
                _rect(slide, 9.55, y + 0.07, 0.12, 0.12, scol)
                _label(slide, 0.6, y, 8.8, 0.38, title, size=10.5, color=OFF_WHITE)
                y += 0.42

    else:  # overview (default)
        _slide_findings(prs, findings[:6], "Top Findings",
                        f"{len(findings)} total data points analysed")
        if len(findings) > 6:
            _slide_findings(prs, findings[6:12], "More Findings")

    # Last slide — Sources
    _slide_sources(prs, findings)

    # End slide — ThreadIntel branding
    _slide_end(prs, topic)

    # ── Save ──────────────────────────────────────────────────────────────────
    if output_path:
        out = Path(output_path)
    else:
        slug = re.sub(r"[^\w]+", "_", topic)[:40]
        out  = ROOT / ".tmp" / f"slides_{slug}_{n.strftime('%Y%m%d_%H%M%S')}.pptx"
    out.parent.mkdir(exist_ok=True)
    prs.save(str(out))

    print(f"\n  PPTX export complete — {len(prs.slides)} slides, {len(findings)} findings")
    print(f"  Saved: {out}\n")
    return str(out)


# ─── Google Slides upload ─────────────────────────────────────────────────────

def export_to_google_slides(
    topic: str,
    findings=None,
    sources: list = None,
    layout: str = "auto",
    share_with: str = None,
) -> str:
    from googleapiclient.discovery import build
    from googleapiclient.http import MediaFileUpload
    from google.oauth2.credentials import Credentials
    from google_auth_oauthlib.flow import InstalledAppFlow
    from google.auth.transport.requests import Request

    DRIVE_SCOPES = ["https://www.googleapis.com/auth/drive.file"]
    token_path   = ROOT / "token_drive.json"
    creds = None
    if token_path.exists():
        creds = Credentials.from_authorized_user_file(str(token_path), DRIVE_SCOPES)
    if not creds or not creds.valid:
        if creds and creds.expired and creds.refresh_token:
            creds.refresh(Request())
        else:
            flow  = InstalledAppFlow.from_client_secrets_file(str(CREDS_FILE), DRIVE_SCOPES)
            creds = flow.run_local_server(port=0)
        token_path.write_text(creds.to_json())

    pptx_path = export_to_pptx(topic, findings=findings, sources=sources, layout=layout)

    drive = build("drive", "v3", credentials=creds)
    n     = datetime.now()
    name  = f"ThreadIntel — {topic[:60]} ({n.strftime(f'%b {n.day}')})"

    uploaded = drive.files().create(
        body={"name": name, "mimeType": "application/vnd.google-apps.presentation"},
        media_body=MediaFileUpload(
            pptx_path,
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            resumable=True,
        ),
        fields="id",
    ).execute()

    file_id = uploaded["id"]
    url = f"https://docs.google.com/presentation/d/{file_id}"

    if share_with:
        drive.permissions().create(
            fileId=file_id,
            body={"type": "user", "role": "writer", "emailAddress": share_with},
            sendNotificationEmail=False,
        ).execute()
        print(f"  Shared with {share_with}")

    print(f"  Uploaded to Google Slides: {url}\n")
    return url


# ─── CLI ──────────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser(
        description="Export ThreadIntel research to PowerPoint / Google Slides"
    )
    parser.add_argument("topic", nargs="?", help="Research topic")
    parser.add_argument("--output",  metavar="FILE", help="Output .pptx path")
    parser.add_argument("--layout",
                        choices=["auto","overview","detailed","comparison","bullets"],
                        default="auto")
    parser.add_argument("--slides",  action="store_true",
                        help="Upload to Google Slides and print URL")
    parser.add_argument("--share",   metavar="EMAIL",
                        help="Share with this email (Google Slides only)")
    args = parser.parse_args()

    if not args.topic:
        args.topic = input("Topic: ").strip()

    if args.slides:
        export_to_google_slides(
            topic=args.topic,
            layout=args.layout,
            share_with=args.share,
        )
    else:
        export_to_pptx(
            topic=args.topic,
            layout=args.layout,
            output_path=args.output,
        )
